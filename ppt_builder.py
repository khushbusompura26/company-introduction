"""
ppt_builder.py — Reference-matched PPT design
Closely follows the Shivalal Agarwala reference presentation:
  · Full-width navy header bar with title inside
  · Alternating White / Light-blue backgrounds
  · Dark navy hero slides for founder / vision / growth
  · Timeline with vertical line + year badges
  · 3-column card grid with dark headers
  · Two-column label:value layout
  · Stats bar at bottom of title slide
  · Consistent footer on every slide
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (from reference PPT) ──────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x37, 0x48)   # body text
GREY   = RGBColor(0x71, 0x80, 0x96)   # secondary text
LGREY  = RGBColor(0xE2, 0xE8, 0xF2)   # divider lines
SBLU   = RGBColor(0x8E, 0xA8, 0xC8)   # stat labels / footer text
CREAM  = RGBColor(0xF5, 0xE6, 0xC0)   # header subtitle text
LBLU   = RGBColor(0xF4, 0xF7, 0xFB)   # light slide background
FTBG   = RGBColor(0x0A, 0x18, 0x30)   # footer bar background
IN     = Inches


# ── Helpers ────────────────────────────────────────────────────

def _adaptive_fs(text, base_fs, ideal_chars=8):
    n = len(str(text))
    if n <= ideal_chars: return base_fs
    return max(8.0, base_fs * (ideal_chars / n) ** 0.65)

def _card_fs(text, base_fs, ideal_chars=55):
    n = len(str(text))
    if n <= ideal_chars: return base_fs
    return max(7.0, base_fs * (ideal_chars / n) ** 0.55)

def _smart_split(point):
    """Split key_point into (header, body) for card/row display."""
    b = str(point).strip()
    if ':' in b:
        hdr, body = b.split(':', 1)
        return hdr.strip()[:60], body.strip()
    sentences = b.split('. ')
    if len(sentences) >= 2 and len(sentences[0]) <= 70:
        return sentences[0].strip(), '. '.join(sentences[1:]).strip()
    words = b.split()
    return ' '.join(words[:5]), b

def _clip(text, n):
    text = str(text)
    return text[:n].rstrip() + '…' if len(text) > n else text

def _founding_year(stats):
    if not stats: return ''
    for k, v in stats.items():
        if any(w in k.lower() for w in ('found','established','estab','incorp','since','year')):
            return str(v)
    return ''


# ── Theme ──────────────────────────────────────────────────────

class Theme:
    def __init__(self, primary_hex='1A2E5A', accent_hex='C9962B'):
        self.primary_hex = primary_hex.upper().lstrip('#')
        self.accent_hex  = accent_hex.upper().lstrip('#')
        self.primary  = self._rgb(self.primary_hex)
        self.accent   = self._rgb(self.accent_hex)
        self.dark_bg  = self._darken(self.primary_hex, 0.60)
        self.darker   = self._darken(self.primary_hex, 0.40)
        self.card_hdr = self._darken(self.primary_hex, 0.65)
        self.chrome   = self._lighten(self.primary_hex, 0.55)
        self.soft     = self._lighten(self.primary_hex, 0.88)

    @staticmethod
    def _rgb(h):
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    @staticmethod
    def _lighten(h, f):
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return RGBColor(min(255,int(r+(255-r)*f)), min(255,int(g+(255-g)*f)), min(255,int(b+(255-b)*f)))
    @staticmethod
    def _darken(h, f):
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return RGBColor(int(r*f), int(g*f), int(b*f))


# ── Low-level drawing ──────────────────────────────────────────

def _rect(sl, l, t, w, h, fill, line=None, lw=Pt(0)):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw or Pt(1)
    else:    sh.line.fill.background()
    return sh

def _rrect(sl, l, t, w, h, fill, line=None, lw=Pt(1.2)):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def _oval(sl, cx, cy, r, fill, line=None, lw=Pt(0)):
    """Draw a circle centred at (cx, cy) with radius r — all in raw EMU/points."""
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def _icon_circle(sl, cx, cy, r, bg, label, fs, fg, fname='Georgia'):
    """Filled circle with centred text (icon / number badge)."""
    _oval(sl, cx, cy, r, bg)
    d = r * 2
    _txt(sl, cx - r, cy - r, d, d, label, fs, fg,
         bold=True, align=PP_ALIGN.CENTER, fname=fname)

def _txt(sl, l, t, w, h, text, fs, color,
         bold=False, italic=False, align=PP_ALIGN.LEFT,
         fname='Calibri', wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = str(text)
    r.font.size = Pt(fs); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = fname
    return tb

def _multi_txt(sl, l, t, w, h, items, fs, color, line_gap=4, fname='Calibri'):
    """Multi-line text box — each item is a paragraph."""
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_gap)
        r = p.add_run(); r.text = '▸  ' + str(item)
        r.font.size = Pt(fs); r.font.color.rgb = color; r.font.name = fname
    return tb

def _bg(sl, color):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = color


# ── Shared components ──────────────────────────────────────────

def _header(sl, t, title, subtitle=''):
    """Full-width header bar — identical to reference PPT."""
    _rect(sl, IN(0),   IN(0),    IN(10), IN(1.15), t.primary)
    _rect(sl, IN(0),   IN(1.15), IN(10), IN(0.07), t.accent)
    _txt(sl, IN(0.4), IN(0.10), IN(9.2), IN(0.65),
         title, 30, WHITE, bold=True, fname='Georgia')
    if subtitle:
        _txt(sl, IN(0.4), IN(0.75), IN(9.2), IN(0.35),
             _clip(subtitle, 130), 11, CREAM, fname='Calibri')

def _footer(sl, t, cn, wu):
    """Footer bar — company | website."""
    _rect(sl, IN(0), IN(5.28), IN(10), IN(0.35), FTBG)
    _txt(sl, IN(0.3), IN(5.29), IN(9.4), IN(0.30),
         f'{cn}  |  {wu}', 8.5, SBLU, fname='Calibri')

def _divider(sl, x, y, w):
    _rect(sl, IN(x), IN(y), IN(w), IN(0.01), LGREY)

def _accent_bar(sl, t, x, y, h, color=None):
    """Thin left accent bar for content panels."""
    _rect(sl, IN(x), IN(y), IN(0.07), IN(h), color or t.primary)

def _section_head(sl, t, x, y, w, text, fs=14):
    _txt(sl, IN(x), IN(y), IN(w), IN(0.40), text,
         fs, t.primary, bold=True, fname='Georgia')

def _lv_row(sl, t, x, y, label, value, lw=1.6, vw=2.6):
    """Label : Value row."""
    _txt(sl, IN(x),         IN(y), IN(lw),  IN(0.27), label, 9.5, t.primary, bold=True, fname='Calibri')
    vfs = _adaptive_fs(str(value), 9.5, ideal_chars=35)
    _txt(sl, IN(x+lw+0.05), IN(y), IN(vw),  IN(0.27), _clip(str(value), 80), vfs, DARK, fname='Calibri')
    _divider(sl, x, y+0.30, lw+vw)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# Dark navy · Big company name · Stats bar at bottom
# ══════════════════════════════════════════════════════════════

def _R_title(prs, layout, t, sd, data):
    sl = prs.slides.add_slide(layout)
    cn      = data.get('company_name', '')
    wu      = data.get('website', '')
    tagline = data.get('tagline', '') or ''
    contact = data.get('contact', {})
    stats   = sd.get('stats') or {}

    _bg(sl, t.primary)

    # ── Decorative background circles ──────────────────────────
    # Large circle — bottom-right quadrant
    _oval(sl, IN(8.8),  IN(5.2),  IN(2.6), t.darker)
    _oval(sl, IN(9.2),  IN(5.0),  IN(1.9), t.card_hdr)
    # Medium circle — top-left
    _oval(sl, IN(0.0),  IN(0.4),  IN(1.8), t.darker)
    _oval(sl, IN(-0.2), IN(0.2),  IN(1.3), t.card_hdr)
    # Small accent dot — top-right
    _oval(sl, IN(9.6),  IN(0.5),  IN(0.5), t.accent)
    _oval(sl, IN(0.5),  IN(5.0),  IN(0.4), t.accent)

    _rect(sl, IN(0), IN(0),    IN(10), IN(0.12), t.darker)
    _rect(sl, IN(0), IN(5.50), IN(10), IN(0.12), t.darker)

    # Est. badge
    found_yr = _founding_year(stats if isinstance(stats, dict) else {})
    if found_yr:
        _rrect(sl, IN(0.5), IN(0.35), IN(1.5), IN(0.38), t.card_hdr)
        _txt(sl, IN(0.5), IN(0.35), IN(1.5), IN(0.38),
             f'Est. {found_yr}', 11, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Calibri')

    # Company name
    _txt(sl, IN(0.5), IN(1.05), IN(9.0), IN(1.1),
         cn, 46, WHITE, bold=True, align=PP_ALIGN.CENTER, fname='Georgia', wrap=True)

    # Tagline (short only — no long description)
    tag = _clip(tagline, 130)
    if tag:
        _txt(sl, IN(0.5), IN(2.25), IN(8.5), IN(0.50), tag, 18, CREAM,
             align=PP_ALIGN.CENTER, fname='Calibri')

    # Industry / short desc
    desc = _clip(sd.get('description', ''), 110)
    if desc:
        _txt(sl, IN(0.5), IN(2.85), IN(8.5), IN(0.40), desc, 13, SBLU,
             italic=True, align=PP_ALIGN.CENTER, fname='Calibri')

    # Contact line
    parts = []
    addr  = contact.get('address', '')
    phone = contact.get('phone', '')
    email = contact.get('email', '')
    if addr:  parts.append(f'📍 {_clip(addr, 35)}')
    if wu:    parts.append(f'🌐 {wu}')
    if phone: parts.append(f'📞 {phone}')
    if email: parts.append(f'✉ {_clip(email, 35)}')
    if parts:
        _txt(sl, IN(0.5), IN(3.42), IN(9.0), IN(0.38),
             '   |   '.join(parts[:4]), 11, SBLU, align=PP_ALIGN.CENTER, fname='Calibri')

    # Stats bar
    if stats and isinstance(stats, dict):
        items = list(stats.items())[:5]
        _rect(sl, IN(0), IN(4.10), IN(10), IN(1.40), t.darker)
        _rect(sl, IN(0), IN(4.10), IN(10), IN(0.06), t.accent)
        ns   = len(items)
        col  = 9.4 / ns
        for i, (lbl, val) in enumerate(items):
            x = 0.3 + i * col
            vfs = _adaptive_fs(str(val), 30, ideal_chars=6)
            _txt(sl, IN(x), IN(4.20), IN(col-0.1), IN(0.65),
                 str(val), vfs, t.accent, bold=True, fname='Georgia')
            _txt(sl, IN(x), IN(4.88), IN(col-0.1), IN(0.35),
                 lbl, 9.5, SBLU, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# LABEL:VALUE TWO-COLUMN
# For: Industry, Overview, Profile, Classification slides
# Header bar · Two panels with label:value rows
# ══════════════════════════════════════════════════════════════

def _R_lv_two_col(prs, layout, t, sd, cn, wu, n, total, bg=WHITE):
    sl = prs.slides.add_slide(layout)
    _bg(sl, bg)
    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])
    stats = sd.get('stats')

    _header(sl, t, title, desc)
    _footer(sl, t, cn, wu)

    CY = 1.35; CH = 3.70
    LW = 4.55; RX = 5.08; RW = 4.55

    mid = max(1, len(pts) // 2)
    left_pts  = pts[:mid]
    right_pts = pts[mid:]

    # Left panel
    _rect(sl, IN(0.28), IN(CY), IN(LW), IN(CH), t.soft if bg == WHITE else WHITE,
          LGREY, Pt(0.5))
    _accent_bar(sl, t, 0.28, CY, CH, t.primary)

    y = CY + 0.15
    for pt in left_pts[:10]:
        if y > CY + CH - 0.25: break
        lbl, val = _smart_split(pt)
        _lv_row(sl, t, 0.48, y, lbl, val, lw=1.55, vw=LW-1.75)
        y += 0.34

    # Right panel
    _rect(sl, IN(RX), IN(CY), IN(RW), IN(CH), t.soft if bg == WHITE else WHITE,
          LGREY, Pt(0.5))
    _accent_bar(sl, t, RX, CY, CH, t.accent)

    # Stats grid on right if available
    if stats and isinstance(stats, dict) and not right_pts:
        _section_head(sl, t, RX+0.2, CY+0.1, RW-0.3, 'Key Metrics', fs=13)
        for i, (k, v) in enumerate(list(stats.items())[:8]):
            ry = CY + 0.55 + i * 0.38
            if ry > CY + CH - 0.25: break
            _lv_row(sl, t, RX+0.2, ry, k, v, lw=1.55, vw=RW-1.75)
    else:
        y = CY + 0.15
        for pt in right_pts[:10]:
            if y > CY + CH - 0.25: break
            lbl, val = _smart_split(pt)
            _lv_row(sl, t, RX+0.2, y, lbl, val, lw=1.55, vw=RW-1.75)
            y += 0.34


# ══════════════════════════════════════════════════════════════
# BULLET TWO-COLUMN
# For general content slides with many points
# Header bar · Two bullet columns
# ══════════════════════════════════════════════════════════════

def _R_two_col(prs, layout, t, sd, cn, wu, n, total, bg=WHITE):
    sl = prs.slides.add_slide(layout)
    _bg(sl, bg)
    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])
    stats = sd.get('stats')

    _header(sl, t, title, desc)
    _footer(sl, t, cn, wu)

    has_stats = bool(stats and isinstance(stats, dict) and stats)
    CY = 1.35; CH = 3.70 if not has_stats else 3.20
    LW = 4.55; RX = 5.08; RW = 4.55

    mid  = max(1, len(pts) // 2)
    lpts = pts[:mid]; rpts = pts[mid:]

    # Left
    _rect(sl, IN(0.28), IN(CY), IN(LW), IN(CH), t.soft if bg == WHITE else WHITE,
          LGREY, Pt(0.5))
    _accent_bar(sl, t, 0.28, CY, CH, t.primary)
    if lpts:
        _multi_txt(sl, IN(0.48), IN(CY+0.15), IN(LW-0.25), IN(CH-0.25),
                   [_clip(p, 120) for p in lpts], 11, DARK)

    # Right
    _rect(sl, IN(RX), IN(CY), IN(RW), IN(CH), t.soft if bg == WHITE else WHITE,
          LGREY, Pt(0.5))
    _accent_bar(sl, t, RX, CY, CH, t.accent)
    if rpts:
        _multi_txt(sl, IN(RX+0.2), IN(CY+0.15), IN(RW-0.25), IN(CH-0.25),
                   [_clip(p, 120) for p in rpts], 11, DARK)

    # Stats pills at bottom
    if has_stats:
        items = list(stats.items())[:5]
        ns = len(items); sw = 9.44 / ns
        sy = CY + CH + 0.15
        for i, (lbl, val) in enumerate(items):
            sx = 0.28 + i * sw
            _rrect(sl, IN(sx), IN(sy), IN(sw-0.05), IN(0.62), t.primary)
            vfs = _adaptive_fs(str(val), 16, ideal_chars=6)
            _txt(sl, IN(sx+0.05), IN(sy+0.04), IN(sw-0.15), IN(0.34),
                 str(val), vfs, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')
            _txt(sl, IN(sx+0.05), IN(sy+0.38), IN(sw-0.15), IN(0.22),
                 lbl, 7.5, WHITE, align=PP_ALIGN.CENTER, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# 3-COLUMN CARD GRID
# Header bar · 3×2 cards with dark navy header bars
# For: Customers, Strengths, Products, Team, etc.
# ══════════════════════════════════════════════════════════════

def _R_cards(prs, layout, t, sd, cn, wu, n, total, bg=WHITE):
    sl = prs.slides.add_slide(layout)
    _bg(sl, bg)
    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])
    stats = sd.get('stats')

    _header(sl, t, title, desc)
    _footer(sl, t, cn, wu)

    has_stats = bool(stats and isinstance(stats, dict) and stats)
    avail_h   = (4.65 if has_stats else 5.18) - 1.40  # content area

    COLS   = 3
    CARD_W = 3.00
    GAP    = (9.44 - COLS * CARD_W) / (COLS - 1)   # ~0.22"
    xs     = [0.28 + i * (CARD_W + GAP) for i in range(COLS)]

    rows   = min((len(pts) + COLS - 1) // COLS, 2)
    rows   = max(rows, 1)
    CARD_H = min((avail_h - 0.18 * (rows - 1)) / rows, 1.85)

    hdr_colors = [t.primary, t.dark_bg, t.card_hdr,
                  t.dark_bg, t.primary, t.card_hdr]

    for i, pt in enumerate(pts[:COLS * rows]):
        row = i // COLS; col = i % COLS
        x = xs[col]; y = 1.40 + row * (CARD_H + 0.18)

        hdr, body = _smart_split(pt)
        body_d    = _clip(body, 160)

        # Card background
        cb = _rrect(sl, IN(x), IN(y), IN(CARD_W), IN(CARD_H),
                    LBLU if bg == WHITE else WHITE)
        cb.line.color.rgb = LGREY; cb.line.width = Pt(0.5)

        # Dark header bar
        _rect(sl, IN(x), IN(y), IN(CARD_W), IN(0.42), hdr_colors[i])

        # Numbered icon circle on the header bar (right-side)
        _icon_circle(sl,
            cx = IN(x + CARD_W - 0.25), cy = IN(y + 0.21),
            r  = IN(0.17), bg = t.accent,
            label = str(i + 1), fs = 8, fg = t.dark_bg)

        # Card title (inside header bar, leave room for circle badge)
        hfs = _card_fs(hdr, 11.0, ideal_chars=35)
        _txt(sl, IN(x+0.12), IN(y+0.07), IN(CARD_W-0.55), IN(0.32),
             hdr, hfs, WHITE, bold=True, fname='Georgia', wrap=True)

        # Card body
        bfs = _card_fs(body_d, 9.5, ideal_chars=80)
        _txt(sl, IN(x+0.12), IN(y+0.48), IN(CARD_W-0.24), IN(CARD_H-0.60),
             body_d, bfs, DARK, fname='Calibri', wrap=True)

    # Stats row
    if has_stats:
        items = list(stats.items())[:5]
        _rect(sl, IN(0.28), IN(4.65), IN(9.44), IN(0.60), t.darker)
        ns = len(items); sw = 9.44 / ns
        for i, (lbl, val) in enumerate(items):
            sx = 0.28 + i * sw
            vfs = _adaptive_fs(str(val), 14, ideal_chars=6)
            _txt(sl, IN(sx+0.05), IN(4.68), IN(sw-0.1), IN(0.30),
                 str(val), vfs, t.accent, bold=True, fname='Georgia')
            _txt(sl, IN(sx+0.05), IN(4.97), IN(sw-0.1), IN(0.22),
                 lbl, 8.5, SBLU, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# TIMELINE
# Header bar · Vertical line · Year badge · Content box
# Exactly 4 entries per slide
# ══════════════════════════════════════════════════════════════

def _R_timeline(prs, layout, t, sd, cn, wu, n, total, bg=LBLU):
    sl = prs.slides.add_slide(layout)
    _bg(sl, bg)
    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])

    _header(sl, t, title, desc)
    _footer(sl, t, cn, wu)

    if not pts: return

    # Vertical line
    _rect(sl, IN(1.60), IN(1.32), IN(0.08), IN(3.72), t.chrome)

    n_items = min(len(pts), 4)
    ih      = 3.72 / n_items
    sy      = 1.32

    badge_colors = [t.primary, t.dark_bg, t.primary, t.dark_bg]

    for i, b in enumerate(pts[:n_items]):
        y = sy + i * ih

        # Parse  "YEAR: Title — Body"  or  "YEAR: Full text"
        if ':' in b:
            yr   = b.split(':', 1)[0].strip()[:12]
            rest = b.split(':', 1)[1].strip()
        else:
            yr   = '—'
            rest = b

        # Split headline from detail
        if ' — ' in rest:
            headline, detail = rest.split(' — ', 1)
        elif '. ' in rest and len(rest.split('. ', 1)[0]) <= 90:
            headline, detail = rest.split('. ', 1)
        else:
            headline = rest[:90]
            detail   = rest[90:] if len(rest) > 90 else ''

        # Year badge
        _rrect(sl, IN(0.28), IN(y+0.06), IN(1.10), IN(0.32), badge_colors[i])
        _txt(sl, IN(0.28), IN(y+0.06), IN(1.10), IN(0.32),
             yr, 9.5, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')

        # Dot on timeline
        _rrect(sl, IN(1.56), IN(y+0.13), IN(0.16), IN(0.16), t.accent)

        # Content box
        ch = ih - 0.12
        cb = _rrect(sl, IN(1.85), IN(y+0.04), IN(7.85), IN(ch),
                    WHITE if bg != WHITE else LBLU)
        cb.line.color.rgb = LGREY; cb.line.width = Pt(0.5)
        _rect(sl, IN(1.85), IN(y+0.04), IN(0.07), IN(ch), t.accent)

        hfs = _card_fs(headline, 10.5, ideal_chars=60)
        _txt(sl, IN(2.05), IN(y+0.07), IN(7.55), IN(0.28),
             headline, hfs, t.primary, bold=True, fname='Calibri')
        if detail:
            dfs = _card_fs(detail, 9.0, ideal_chars=120)
            _txt(sl, IN(2.05), IN(y+0.38), IN(7.55), IN(ch-0.40),
                 _clip(detail, 200), dfs, DARK, fname='Calibri', wrap=True)


# ══════════════════════════════════════════════════════════════
# PROCESS / OPERATIONS
# Header bar · Numbered steps (left) · Locations/info (right)
# ══════════════════════════════════════════════════════════════

def _R_process(prs, layout, t, sd, cn, wu, n, total, bg=WHITE):
    sl = prs.slides.add_slide(layout)
    _bg(sl, bg)
    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])

    _header(sl, t, title, desc)
    _footer(sl, t, cn, wu)

    CY = 1.35; CH = 3.70
    LW = 5.55; RX = 6.05; RW = 3.68

    # Left panel
    _rect(sl, IN(0.28), IN(CY), IN(LW), IN(CH),
          t.soft if bg == WHITE else WHITE, LGREY, Pt(0.5))
    _accent_bar(sl, t, 0.28, CY, CH, t.primary)

    parts = title.split('&')
    left_head  = parts[0].strip() if parts else 'Process'
    right_head = parts[-1].strip() if len(parts) > 1 else 'Operations'
    _section_head(sl, t, 0.48, CY+0.08, LW-0.3, left_head, fs=13)

    step_h = 0.54
    for i, pt in enumerate(pts[:6]):
        sy2 = CY + 0.55 + i * step_h
        if sy2 + 0.4 > CY + CH: break
        hdr, body = _smart_split(pt)
        _rrect(sl, IN(0.42), IN(sy2+0.03), IN(0.32), IN(0.32), t.primary)
        _txt(sl, IN(0.42), IN(sy2+0.03), IN(0.32), IN(0.32),
             f'{i+1:02d}', 8, WHITE, bold=True, align=PP_ALIGN.CENTER, fname='Calibri')
        _txt(sl, IN(0.88), IN(sy2+0.01), IN(LW-0.65), IN(0.25),
             _clip(hdr, 45), 9.5, t.primary, bold=True, fname='Calibri')
        bfs = _card_fs(body, 8.5, ideal_chars=80)
        _txt(sl, IN(0.88), IN(sy2+0.27), IN(LW-0.65), IN(0.23),
             _clip(body, 130), bfs, DARK, fname='Calibri', wrap=True)

    # Right panel
    _rect(sl, IN(RX), IN(CY), IN(RW), IN(CH),
          t.soft if bg == WHITE else WHITE, LGREY, Pt(0.5))
    _accent_bar(sl, t, RX, CY, CH, t.accent)
    _section_head(sl, t, RX+0.2, CY+0.08, RW-0.3, right_head, fs=13)

    right_pts = pts[6:] if len(pts) > 6 else pts[4:] if len(pts) > 4 else pts[3:]
    for i, pt in enumerate(right_pts[:5]):
        ry = CY + 0.55 + i * 0.66
        if ry + 0.4 > CY + CH: break
        hdr, body = _smart_split(pt)
        _txt(sl, IN(RX+0.2), IN(ry),       IN(RW-0.3), IN(0.28),
             f'▸  {_clip(hdr, 40)}', 10.0, t.primary, bold=True, fname='Calibri')
        bfs = _card_fs(body, 8.5, ideal_chars=70)
        _txt(sl, IN(RX+0.2), IN(ry+0.30), IN(RW-0.3), IN(0.28),
             _clip(body, 100), bfs, DARK, fname='Calibri', wrap=True)
        _divider(sl, RX+0.2, ry+0.62, RW-0.3)


# ══════════════════════════════════════════════════════════════
# DARK HERO SLIDE
# Dark navy bg · Portrait area (left) · Rich content (right)
# For: Founder, Vision/Mission, Growth, Digital, Testimonials
# ══════════════════════════════════════════════════════════════

def _R_dark(prs, layout, t, sd, cn, wu, n, total):
    sl = prs.slides.add_slide(layout)
    _bg(sl, t.dark_bg)

    title = sd.get('title', '')
    desc  = sd.get('description', '')
    pts   = sd.get('key_points', [])
    stats = sd.get('stats')

    # ── Decorative background circles ──────────────────────────
    _oval(sl, IN(9.5),  IN(5.0),  IN(2.2), t.darker)
    _oval(sl, IN(10.0), IN(5.4),  IN(1.5), t.card_hdr)
    _oval(sl, IN(0.2),  IN(5.2),  IN(1.1), t.darker)
    _oval(sl, IN(9.8),  IN(0.2),  IN(0.6), t.accent)

    # Thin accent lines at top
    _rect(sl, IN(0), IN(0),    IN(10), IN(0.10), t.darker)
    _rect(sl, IN(0), IN(0.10), IN(10), IN(0.04), t.accent)

    # Slide counter
    _txt(sl, IN(9.2), IN(0.18), IN(0.75), IN(0.28),
         f'{n}/{total}', 8, t.chrome, align=PP_ALIGN.RIGHT, fname='Calibri')

    # Left portrait panel
    _rrect(sl, IN(0.45), IN(0.70), IN(2.80), IN(2.80), t.card_hdr)
    initials = ''.join(w[0].upper() for w in title.replace('&','').split()[:2])
    _txt(sl, IN(0.45), IN(0.70), IN(2.80), IN(2.80),
         initials or 'CO', 52, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')

    # Name / role below portrait
    name_text = pts[0] if pts else title
    name_text = _smart_split(name_text)[0] if ':' in name_text else _clip(name_text, 40)
    _txt(sl, IN(0.45), IN(3.60), IN(2.80), IN(0.35), name_text, 12, CREAM,
         align=PP_ALIGN.CENTER, fname='Calibri')
    _txt(sl, IN(0.45), IN(3.95), IN(2.80), IN(0.30), cn, 10, SBLU,
         align=PP_ALIGN.CENTER, fname='Calibri')

    # Divider line
    _rect(sl, IN(3.60), IN(0.60), IN(0.06), IN(4.50), t.chrome)

    # Right — title, subtitle, content
    _txt(sl, IN(3.90), IN(0.20), IN(5.90), IN(0.55),
         title, 30, WHITE, bold=True, fname='Georgia', wrap=True)
    if desc:
        _txt(sl, IN(3.90), IN(0.78), IN(5.90), IN(0.38),
             _clip(desc, 130), 12, CREAM, italic=True, fname='Calibri')
    _rect(sl, IN(3.90), IN(1.25), IN(5.90), IN(0.04), t.accent)

    # Main body text from remaining points
    body_pts = pts[1:] if len(pts) > 1 else pts
    if body_pts:
        body_text = ' '.join(_clip(p, 120) for p in body_pts[:4])
        _txt(sl, IN(3.90), IN(1.38), IN(5.85), IN(2.85),
             body_text, 11.5, SBLU, fname='Calibri', wrap=True)

    # Stats bar at bottom right
    if stats and isinstance(stats, dict):
        items = list(stats.items())[:5]
        _rect(sl, IN(3.90), IN(4.45), IN(5.90), IN(0.75), t.darker)
        sw = 5.90 / len(items)
        for i, (lbl, val) in enumerate(items):
            sx = 3.90 + i * sw
            vfs = _adaptive_fs(str(val), 11, ideal_chars=6)
            _txt(sl, IN(sx+0.05), IN(4.52), IN(sw-0.1), IN(0.28),
                 str(val), vfs, t.accent, bold=True, fname='Georgia')
            _txt(sl, IN(sx+0.05), IN(4.80), IN(sw-0.1), IN(0.22),
                 lbl, 8.5, SBLU, fname='Calibri')

    # Footer
    _rect(sl, IN(0), IN(5.28), IN(10), IN(0.35), t.darker)
    _txt(sl, IN(0.3), IN(5.29), IN(9.4), IN(0.30),
         f'{cn}  |  {wu}', 8.5, GREY, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# QUOTE / IMPACT SLIDE
# Dark bg · Large statement text · 3 supporting pillars
# For: Vision, Mission, Values, CSR, Awards
# ══════════════════════════════════════════════════════════════

def _R_quote(prs, layout, t, sd, cn, wu, n, total):
    sl = prs.slides.add_slide(layout)
    _bg(sl, t.dark_bg)

    title = sd.get('title', '')
    desc  = sd.get('description', '') or ''
    pts   = sd.get('key_points', [])

    # ── Decorative circles ──────────────────────────────────────
    _oval(sl, IN(-0.5), IN(5.6),  IN(2.5), t.darker)
    _oval(sl, IN(10.5), IN(-0.3), IN(2.2), t.darker)
    _oval(sl, IN(10.2), IN(0.0),  IN(1.4), t.card_hdr)
    _oval(sl, IN(0.2),  IN(0.0),  IN(0.5), t.accent)
    _oval(sl, IN(9.8),  IN(5.3),  IN(0.4), t.accent)

    # ── Accent lines ────────────────────────────────────────────
    _rect(sl, IN(0), IN(0),    IN(10), IN(0.10), t.darker)
    _rect(sl, IN(0), IN(0.10), IN(10), IN(0.04), t.accent)

    # ── Slide counter ───────────────────────────────────────────
    _txt(sl, IN(9.2), IN(0.18), IN(0.75), IN(0.28),
         f'{n}/{total}', 8, t.chrome, align=PP_ALIGN.RIGHT, fname='Calibri')

    # ── Title bar ───────────────────────────────────────────────
    _txt(sl, IN(0.5), IN(0.22), IN(9.0), IN(0.45),
         title, 22, WHITE, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')
    _rect(sl, IN(4.0), IN(0.72), IN(2.0), IN(0.05), t.accent)

    # ── Main statement quote (description) ──────────────────────
    if desc:
        statement = _clip(desc, 200)
        # Large decorative quote marks
        _txt(sl, IN(0.3), IN(0.80), IN(0.6), IN(0.55), '“', 40, t.accent,
             bold=True, fname='Georgia')
        _txt(sl, IN(0.7), IN(0.92), IN(8.6), IN(0.85),
             statement, 15, CREAM, italic=True,
             align=PP_ALIGN.CENTER, fname='Calibri', wrap=True)
        _txt(sl, IN(9.1), IN(1.40), IN(0.6), IN(0.55), '”', 40, t.accent,
             bold=True, fname='Georgia', align=PP_ALIGN.RIGHT)

    # ── Three pillar boxes ───────────────────────────────────────
    pillar_pts = pts[:3] if len(pts) >= 3 else (pts + [''] * 3)[:3]
    pillar_colors = [t.primary, t.dark_bg, t.primary]
    border_colors = [t.accent,  t.chrome,  t.accent]
    icon_labels   = ['01', '02', '03']

    PW = 2.90; PH = 1.85; GAP = 0.28; PY = 1.98
    xs = [0.28 + i * (PW + GAP) for i in range(3)]

    for i, (pt, px) in enumerate(zip(pillar_pts, xs)):
        hdr, body = _smart_split(pt) if pt else ('—', '')

        # Pillar card
        cb = _rrect(sl, IN(px), IN(PY), IN(PW), IN(PH), pillar_colors[i],
                    border_colors[i], lw=Pt(1.2))

        # Top accent band
        _rect(sl, IN(px), IN(PY), IN(PW), IN(0.07), border_colors[i])

        # Icon circle at top-centre of card
        _icon_circle(sl,
            cx = IN(px + PW / 2), cy = IN(PY + 0.38),
            r  = IN(0.26), bg = t.accent,
            label = icon_labels[i], fs = 9, fg = t.dark_bg)

        # Pillar title
        hfs = _card_fs(hdr, 11.5, ideal_chars=30)
        _txt(sl, IN(px+0.12), IN(PY+0.70), IN(PW-0.24), IN(0.32),
             hdr, hfs, WHITE, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')

        # Pillar body
        if body:
            bfs = _card_fs(body, 9.0, ideal_chars=80)
            _txt(sl, IN(px+0.14), IN(PY+1.04), IN(PW-0.28), IN(PH-1.12),
                 _clip(body, 130), bfs, SBLU, align=PP_ALIGN.CENTER,
                 fname='Calibri', wrap=True)

    # ── Footer ───────────────────────────────────────────────────
    _rect(sl, IN(0), IN(5.28), IN(10), IN(0.35), t.darker)
    _txt(sl, IN(0.3), IN(5.29), IN(9.4), IN(0.30),
         f'{cn}  |  {wu}', 8.5, GREY, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# CLOSING SLIDE
# Dark bg · Thank you · 2×2 contact cards
# ══════════════════════════════════════════════════════════════

def _R_closing(prs, layout, t, sd, data):
    sl  = prs.slides.add_slide(layout)
    cn  = data.get('company_name', '')
    wu  = data.get('website', '')
    con = data.get('contact', {})

    _bg(sl, t.primary)

    # ── Decorative background circles ──────────────────────────
    _oval(sl, IN(9.8),  IN(5.2),  IN(2.4), t.darker)
    _oval(sl, IN(10.2), IN(5.5),  IN(1.6), t.card_hdr)
    _oval(sl, IN(-0.2), IN(0.3),  IN(1.6), t.darker)
    _oval(sl, IN(0.0),  IN(5.0),  IN(0.8), t.accent)
    _oval(sl, IN(9.8),  IN(0.4),  IN(0.5), t.accent)

    _rect(sl, IN(0), IN(0),    IN(10), IN(0.10), t.darker)
    _rect(sl, IN(0), IN(0.10), IN(10), IN(0.04), t.accent)

    # Thank You
    _txt(sl, IN(0.5), IN(0.40), IN(9.0), IN(1.10),
         'Thank You', 52, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')
    _rect(sl, IN(3.5), IN(1.55), IN(3.0), IN(0.06), t.accent)

    # Closing message
    desc = sd.get('description', 'Thank you for your time.')
    _txt(sl, IN(0.8), IN(1.72), IN(8.4), IN(0.50),
         _clip(desc, 160), 14, CREAM, italic=True, align=PP_ALIGN.CENTER,
         fname='Calibri', wrap=True)

    # 2×2 contact cards
    cards = [
        ('🌐 Website',  wu or '—'),
        ('📍 Address',  con.get('address', '—')),
        ('📞 Phone',    con.get('phone',   '—')),
        ('✉️ Email',    con.get('email',   '—')),
    ]
    for i, (lbl, val) in enumerate(cards):
        x = 0.25 if i < 2 else 5.12
        y = 2.55 + (i % 2) * 0.82
        cc = _rrect(sl, IN(x), IN(y), IN(4.63), IN(0.70),
                    t.dark_bg, t.accent, lw=Pt(0.8))
        _rect(sl, IN(x), IN(y), IN(4.63), IN(0.07),
              t.accent if i % 2 == 0 else t.chrome)
        _txt(sl, IN(x+0.12), IN(y+0.10), IN(4.4), IN(0.26),
             lbl, 9.5, t.accent, bold=True, fname='Calibri')
        vfs = _card_fs(str(val), 11.0, ideal_chars=50)
        _txt(sl, IN(x+0.12), IN(y+0.36), IN(4.4), IN(0.28),
             _clip(str(val), 60), vfs, WHITE, fname='Calibri')

    # Company + website
    _txt(sl, IN(0.5), IN(4.38), IN(9.0), IN(0.25),
         cn, 12, t.accent, bold=True, align=PP_ALIGN.CENTER, fname='Georgia')
    _rect(sl, IN(0), IN(5.28), IN(10), IN(0.35), t.darker)
    _txt(sl, IN(0.3), IN(5.29), IN(9.4), IN(0.30),
         wu, 10, SBLU, italic=True, align=PP_ALIGN.CENTER, fname='Calibri')


# ══════════════════════════════════════════════════════════════
# LAYOUT PICKER
# ══════════════════════════════════════════════════════════════

def _pick_layout(sd, idx, total):
    n     = sd.get('slide_number', idx + 1)
    title = sd.get('title', '').lower()
    pts   = sd.get('key_points', [])

    if n == 1:      return 'title'
    if n == total:  return 'closing'

    # Quote / impact slides — Vision, Mission, Values, CSR, Awards
    quote_kw = ('vision', 'mission', 'value', 'csr', 'award', 'recogni',
                'purpose', 'philosophy', 'pledge', 'commitment')
    if any(k in title for k in quote_kw):
        return 'quote'

    # Dark hero slides
    dark_kw = ('founder', 'legacy', 'eminent', 'growth', 'scale',
                'digital product', 'toppers', 'testimonial',
                'roadmap', 'future')
    if any(k in title for k in dark_kw):
        return 'dark'

    # Timeline
    timeline_kw = ('history', 'timeline', 'journey', 'milestone', 'evolution')
    if any(k in title for k in timeline_kw):
        return 'timeline'

    # Process / operations
    process_kw = ('operation', 'process', 'manufactur', 'production', 'wareho',
                  'logistic', 'step', 'phase', 'delivery', 'supply')
    if any(k in title for k in process_kw):
        return 'process'

    # Card grid
    cards_kw = ('strength', 'usp', 'unique', 'product', 'service', 'offering',
                'customer', 'client', 'audience', 'target', 'director', 'team',
                'leadership', 'network', 'dealer', 'geographic', 'reach',
                'brand', 'marketing', 'digital presence', 'distribution')
    if any(k in title for k in cards_kw) and len(pts) >= 3:
        return 'cards'

    # Label:value two-col (data-heavy classification slides)
    lv_kw = ('industry', 'business type', 'overview', 'profile',
              'classification', 'about', 'introduction')
    if any(k in title for k in lv_kw):
        return 'lv_two_col'

    # Two-col bullets for 6+ points
    if len(pts) >= 6:
        return 'two_col'

    return 'lv_two_col'


# ══════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ══════════════════════════════════════════════════════════════

def build_presentation(data, primary_hex='1A2E5A', accent_hex='C9962B'):
    t = Theme(primary_hex, accent_hex)

    prs = Presentation()
    prs.slide_width  = IN(10)
    prs.slide_height = IN(5.625)
    blank = prs.slide_layouts[6]

    cn     = data.get('company_name', 'Company')
    wu     = data.get('website', '')
    slides = data.get('slides', [])
    total  = len(slides)

    def bg(idx):
        """Alternate White / Light-blue for regular slides."""
        return WHITE if idx % 2 == 0 else LBLU

    for idx, sd in enumerate(slides):
        n      = sd.get('slide_number', idx + 1)
        layout = _pick_layout(sd, idx, total)

        if layout == 'title':
            _R_title(prs, blank, t, sd, data)
        elif layout == 'closing':
            _R_closing(prs, blank, t, sd, data)
        elif layout == 'quote':
            _R_quote(prs, blank, t, sd, cn, wu, n, total)
        elif layout == 'dark':
            _R_dark(prs, blank, t, sd, cn, wu, n, total)
        elif layout == 'timeline':
            _R_timeline(prs, blank, t, sd, cn, wu, n, total, bg=bg(idx))
        elif layout == 'process':
            _R_process(prs, blank, t, sd, cn, wu, n, total, bg=bg(idx))
        elif layout == 'cards':
            _R_cards(prs, blank, t, sd, cn, wu, n, total, bg=bg(idx))
        elif layout == 'lv_two_col':
            _R_lv_two_col(prs, blank, t, sd, cn, wu, n, total, bg=bg(idx))
        else:  # two_col
            _R_two_col(prs, blank, t, sd, cn, wu, n, total, bg=bg(idx))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), 'reference'
